from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTFILE = "microsoft-foundry-public-overview-15slides.pptx"

BG_DARK = RGBColor(30, 39, 97)
BG_LIGHT = RGBColor(245, 247, 252)
ACCENT = RGBColor(42, 122, 228)
TEXT_DARK = RGBColor(35, 43, 63)
TEXT_MUTED = RGBColor(96, 109, 133)
WHITE = RGBColor(255, 255, 255)


def add_title(slide, text, subtitle=None, dark=False):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = WHITE if dark else TEXT_DARK
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(8.8), Inches(0.5))
        p2 = sub_box.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(15)
        p2.font.color.rgb = WHITE if dark else TEXT_MUTED


def add_footer_source(slide, text, dark=False):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(5.15), Inches(8.8), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(220, 227, 244) if dark else TEXT_MUTED


def add_card(slide, x, y, w, h, title, body, color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(220, 225, 235)
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(16)
    p1.font.color.rgb = TEXT_DARK
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = TEXT_MUTED


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    add_title(slide, "Microsoft Foundry 公开资料解读", "面向企业 AI 应用的统一平台（15 页）", dark=True)
    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.1), Inches(3.2), Inches(0.7))
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT
    tag.line.fill.background()
    tag.text_frame.text = "来源：Microsoft Learn 官方文档"
    tag.text_frame.paragraphs[0].font.size = Pt(14)
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].font.color.rgb = WHITE
    add_footer_source(slide, "https://learn.microsoft.com/azure/ai-foundry/what-is-azure-ai-foundry", dark=True)

    # 2. What is Foundry
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "1) 什么是 Microsoft Foundry", "统一管理模型、智能体与工具的企业级 AI 平台")
    add_card(slide, 0.6, 1.4, 4.2, 1.5, "平台定位", "PaaS 形态，聚焦应用开发与运维治理，减少基础设施管理复杂度。")
    add_card(slide, 5.1, 1.4, 4.3, 1.5, "统一控制面", "将 RBAC、网络、策略与资产管理汇聚到单一资源模型。")
    add_card(slide, 0.6, 3.1, 8.8, 1.6, "核心价值", "以“项目端点 + API/SDK + 门户”统一体验支持构建、评估、监控与规模化落地。", color=RGBColor(235, 243, 255))
    add_footer_source(slide, "来源：What is Azure AI Foundry")

    # 3. Evolution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "2) 平台演进：从分散到统一", "官方文档给出的关键迁移方向")
    headers = ["维度", "过去", "现在"]
    rows = [
        ["品牌", "Azure AI Studio / Foundry", "Microsoft Foundry"],
        ["资源模型", "Hub + 多资源", "Foundry 资源 + 项目"],
        ["Agent API", "Assistants API", "Responses API (v1)"],
        ["SDK与端点", "多端点多包", "统一项目客户端 + 项目端点"],
    ]
    x0, y0 = 0.7, 1.5
    col_w = [1.4, 3.2, 4.1]
    for i, h in enumerate(headers):
        add_card(slide, x0 + sum(col_w[:i]), y0, col_w[i], 0.55, h, "", color=RGBColor(220, 234, 255))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            add_card(slide, x0 + sum(col_w[:c]), y0 + 0.62 + r * 0.85, col_w[c], 0.78, val, "", color=WHITE)
    add_footer_source(slide, "来源：Evolution of Foundry（官方迁移对照）")

    # 4. Audience
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "3) 适用人群与典型价值", "官方定义的三类核心用户")
    add_card(slide, 0.6, 1.5, 2.8, 2.8, "应用开发者", "构建基于模型与智能体的产品，快速集成工具与知识能力。")
    add_card(slide, 3.6, 1.5, 2.8, 2.8, "ML 工程师 / 数据科学家", "进行微调、评估与模型部署运营，提升效果与稳定性。")
    add_card(slide, 6.6, 1.5, 2.8, 2.8, "IT 管理员 / 平台工程师", "执行治理、访问控制、策略与合规审计。")
    add_footer_source(slide, "来源：Who is Foundry for?")

    # 5. Build capabilities
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "4) 构建能力：智能体与工具体系", "面向 agentic 应用的官方能力清单")
    add_card(slide, 0.6, 1.4, 2.75, 1.3, "多智能体编排", "支持协作式工作流与复杂任务拆解。")
    add_card(slide, 3.55, 1.4, 2.75, 1.3, "工具目录", "对接 1,400+ 工具，覆盖公有与私有源。")
    add_card(slide, 6.5, 1.4, 2.9, 1.3, "Memory", "跨会话保留上下文，降低重复输入成本。")
    add_card(slide, 0.6, 3.0, 4.3, 1.5, "Foundry IQ", "连接企业/网页知识，输出可溯源引用结果。")
    add_card(slide, 5.1, 3.0, 4.3, 1.5, "发布能力", "发布到 Microsoft 365、Teams、BizChat 或容器部署。")
    add_footer_source(slide, "来源：Key capabilities - Build agents")

    # 6. Operate and govern
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "5) 运维与治理：企业落地的关键", "性能、可观测、权限与策略统一管理")
    add_card(slide, 0.7, 1.4, 4.2, 1.2, "实时可观测", "持续评估与监控，定位模型/智能体质量与风险。", color=RGBColor(235, 243, 255))
    add_card(slide, 5.1, 1.4, 4.2, 1.2, "资产集中运营", "统一管理跨云注册的模型、智能体与工具。", color=RGBColor(235, 243, 255))
    add_card(slide, 0.7, 2.9, 8.6, 1.4, "企业级控制面", "支持 MCP/A2A 认证、AI Gateway 与 Azure Policy 集成，实现规模治理。")
    add_footer_source(slide, "来源：Key capabilities - Operate and govern")

    # 7. SDK choice
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "6) SDK/API 选型建议", "官方文档：按场景选择最小复杂度路径")
    add_card(slide, 0.6, 1.4, 2.1, 2.7, "Foundry SDK", "项目端点统一接入\n适合 Agents / Evaluations / Foundry 特性")
    add_card(slide, 2.9, 1.4, 2.1, 2.7, "OpenAI SDK", "OpenAI API 兼容优先\n适合 Chat Completions 场景")
    add_card(slide, 5.2, 1.4, 2.1, 2.7, "Tools SDK", "Vision / Speech / Safety 等专项能力")
    add_card(slide, 7.5, 1.4, 1.9, 2.7, "Agent Framework", "多智能体编排代码化\n云无关设计")
    add_footer_source(slide, "来源：SDK Overview（Foundry endpoint vs OpenAI endpoint）")

    # 8. Workflow patterns
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "7) Workflow 编排模式（官方）", "在可视化设计器中构建可重复流程")
    add_card(slide, 0.6, 1.5, 2.8, 2.7, "Human in the loop", "引入审批/澄清节点，满足高风险业务的人机协同。")
    add_card(slide, 3.6, 1.5, 2.8, 2.7, "Sequential", "按固定顺序串联多个 Agent，适合流水线处理。")
    add_card(slide, 6.6, 1.5, 2.8, 2.7, "Group chat", "根据上下文动态切换 Agent，适合专家分工与升级。")
    add_footer_source(slide, "来源：Agents Workflow 概念文档")

    # 9. Cost + region with chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "8) 成本与区域规划", "先估算、再验证、后规模化")
    data = CategoryChartData()
    data.categories = ["模型推理", "工具服务", "微调/托管", "治理与观测"]
    data.add_series("成本构成示例", (45, 25, 20, 10))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.7), Inches(1.45), Inches(4.2), Inches(3.1), data).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    add_card(slide, 5.1, 1.5, 4.2, 1.3, "成本管理关键动作", "1) 定价计算器估算\n2) 小流量压测校准\n3) Cost Management 按 Meter 对账")
    add_card(slide, 5.1, 3.0, 4.2, 1.3, "区域规划关键动作", "先选候选区域，再核验模型/工具可用性与配额。\n官方文档当前列出 30+ Foundry 项目区域。")
    add_footer_source(slide, "来源：Manage costs / Region support")

    # 10. API & endpoint strategy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "9) 端点与 API 策略", "同一 Foundry 资源可支持不同开发路径")
    add_card(slide, 0.6, 1.4, 4.2, 1.6, "项目端点（推荐）", "https://<resource>.services.ai.azure.com/api/projects/<project>\n适合使用 Foundry SDK 与项目级能力。", color=RGBColor(235, 243, 255))
    add_card(slide, 5.1, 1.4, 4.3, 1.6, "OpenAI 兼容端点", "https://<resource>.openai.azure.com/openai/v1\n适合已有 OpenAI SDK 生态与兼容性诉求。", color=RGBColor(235, 243, 255))
    add_card(slide, 0.6, 3.3, 8.8, 1.3, "选型原则", "优先按“治理需求 + 现有技术栈 + 迁移成本”决策，避免多端点并行导致复杂度上升。")
    add_footer_source(slide, "来源：SDK Overview")

    # 11. RBAC and governance roles
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "10) 权限与治理基线", "RBAC 角色分层是生产落地前置条件")
    add_card(slide, 0.6, 1.5, 2.8, 2.7, "Azure AI User", "面向开发与使用，适合最小权限实践。")
    add_card(slide, 3.6, 1.5, 2.8, 2.7, "AI Project Manager", "管理 Foundry 项目与协作流程。")
    add_card(slide, 6.6, 1.5, 2.8, 2.7, "Owner / Contributor", "订阅级资源管理与策略落地。")
    add_footer_source(slide, "来源：SDK Overview / Manage costs / RBAC guidance")

    # 12. Quota and regional rollout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "11) 配额与区域扩展策略", "先验证区域可用性，再按配额推进发布")
    add_card(slide, 0.7, 1.4, 4.2, 1.3, "区域验证步骤", "1) 选候选区域\n2) 核验模型/工具\n3) 检查配额与容量", color=RGBColor(235, 243, 255))
    add_card(slide, 5.1, 1.4, 4.2, 1.3, "配额入口", "Foundry 门户 Operate > Quota，可查看模型与区域可用情况。", color=RGBColor(235, 243, 255))
    add_card(slide, 0.7, 3.0, 8.6, 1.5, "发布建议", "采用“单区域试点 → 双区域容灾 → 多区域复制”路径，同时建立容量阈值告警。")
    add_footer_source(slide, "来源：Region support / Control plane guidance")

    # 13. Sovereign cloud considerations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "12) 主权云与合规场景", "公共云与 Azure Government 的能力差异需要前置评估")
    add_card(slide, 0.6, 1.5, 4.3, 2.7, "Azure Government", "门户：ai.azure.us\n区域：US Gov Arizona / Virginia\n适用：政府与合规敏感行业")
    add_card(slide, 5.1, 1.5, 4.3, 2.7, "关键限制", "部分功能在主权云不可用（如部分 playground、fine-tuning 等），需按官方清单核验。")
    add_footer_source(slide, "来源：Region support（Foundry in sovereign clouds）")

    # 14. Migration checklist
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "13) 从 Azure OpenAI 到 Foundry 的迁移清单", "在保持业务连续性的前提下升级平台能力")
    add_card(slide, 0.6, 1.4, 8.8, 1.1, "步骤 1：资源映射", "确认现有模型、端点、密钥与项目结构的映射关系。", color=RGBColor(235, 243, 255))
    add_card(slide, 0.6, 2.7, 8.8, 1.1, "步骤 2：接口改造", "按 Responses API 与新术语（Conversations/Items/Responses）完成 SDK 调整。", color=RGBColor(235, 243, 255))
    add_card(slide, 0.6, 4.0, 8.8, 1.1, "步骤 3：治理补强", "补齐监控、评估、预算与 RBAC 规则后再切生产流量。", color=RGBColor(235, 243, 255))
    add_footer_source(slide, "来源：What is Foundry（Evolution / Upgrade guidance）")

    # 15. Roadmap + sources
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(22, 28, 44)
    bg.line.fill.background()
    add_title(slide, "14) 落地路线图与官方资料", "建议从小场景试点，逐步扩展到生产", dark=True)
    add_card(slide, 0.6, 1.5, 8.8, 1.5, "90 天落地建议（示例）", "阶段1：搭建 Foundry 项目与最小可用 Agent\n阶段2：接入企业知识与工具，建立评估指标\n阶段3：上线监控、预算告警与治理策略", color=RGBColor(233, 241, 255))
    refs = slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(8.8), Inches(1.9))
    tf = refs.text_frame
    tf.word_wrap = True
    entries = [
        "官方文档入口：https://learn.microsoft.com/azure/ai-foundry/",
        "What is Foundry：https://learn.microsoft.com/en-us/azure/ai-foundry/what-is-azure-ai-foundry",
        "SDK 概览：https://learn.microsoft.com/en-us/azure/ai-foundry/how-to/develop/sdk-overview",
        "Workflow：https://learn.microsoft.com/en-us/azure/ai-foundry/agents/concepts/workflow",
        "成本管理：https://learn.microsoft.com/en-us/azure/ai-foundry/concepts/manage-costs",
        "区域支持：https://learn.microsoft.com/en-us/azure/ai-foundry/reference/region-support",
    ]
    for i, e in enumerate(entries):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = e
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE

    prs.save(OUTFILE)


if __name__ == "__main__":
    build()
